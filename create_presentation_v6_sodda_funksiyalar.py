#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot - Sodda Funksiyalar Prezentatsiyasi
Yaratuvchi: AI Assistant
Sana: 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_simple_functions_presentation():
    """Botning barcha funksiyalarini sodda va tushunarli holda ko'rsatuvchi prezentatsiya yaratish"""
    
    # Yangi prezentatsiya yaratish
    prs = Presentation()
    
    # Slide 1: Sarlavha
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "🎬 Centris Bot"
    subtitle1.text = "Barcha Funksiyalar - Sodda va Tushunarli"
    
    # Slide 2: Bot Nima?
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "🤖 Bot Nima?"
    content2.text = """• Telegram bot
• Centris Towers va Golden Lake loyihalari uchun
• Videolarni avtomatik yuborish
• Foydalanuvchilarga video taqsimlash
• Admin boshqaruvi"""
    
    # Slide 3: Asosiy Funksiyalar
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "⚡ Asosiy Funksiyalar"
    content3.text = """• Video yuborish
• Vaqtni belgilash
• Guruh boshqaruvi
• Foydalanuvchi ro'yxatdan o'tkazish
• Admin huquqlari
• Xavfsizlik tizimi"""
    
    # Slide 4: Foydalanuvchi Buyruqlari
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "👤 Foydalanuvchi Buyruqlari"
    content4.text = """• /start - Botni ishga tushirish
• /help - Yordam olish
• /menu - Asosiy menyu
• /projects - Loyihalarni ko'rish
• /videos - Videolarni ko'rish
• /settings - Sozlamalarni o'zgartirish"""
    
    # Slide 5: Video Funksiyalari
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "🎥 Video Funksiyalari"
    content5.text = """• Avtomatik video yuborish
• Vaqtni o'zi belgilash
• Keyingi ko'rilmagan videoni olish
• Video progress kuzatish
• Video ro'yxatini ko'rish
• Video ma'lumotlarini ko'rish"""
    
    # Slide 6: Guruh Boshqaruvi
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "👥 Guruh Boshqaruvi"
    content6.text = """• Guruh qo'shish
• Guruh o'chirish
• Guruh nomini o'zgartirish
• Guruh sozlamalarini ko'rish
• Guruh videolarini yuborish
• Guruh a'zolarini boshqarish"""
    
    # Slide 7: Admin Buyruqlari
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "🔐 Admin Buyruqlari"
    content7.text = """• /set_group_video - Guruh video sozlamalari
• /add_season - Yangi mavsum qo'shish
• /add_video - Yangi video qo'shish
• /send_video - Videoni yuborish
• /user_management - Foydalanuvchilarni boshqarish
• /statistics - Statistikalarni ko'rish"""
    
    # Slide 8: Vaqt Boshqaruvi
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "⏰ Vaqt Boshqaruvi"
    content8.text = """• Default vaqtlar: 08:00 va 20:00
• O'z vaqtini belgilash
• 5 ta tayyor vaqt: 09:00, 12:00, 15:00, 18:00, 21:00
• Har kuni avtomatik video yuborish
• Vaqtni o'zgartirish
• Scheduler boshqaruvi"""
    
    # Slide 9: Xavfsizlik Tizimi
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "🔒 Xavfsizlik Tizimi"
    content9.text = """• Foydalanuvchi ro'yxatdan o'tkazish
• Admin huquqlari
• Super admin huquqlari
• Guruh whitelist
• Xavfsizlik tekshiruvlari
• Ruxsat berish/olish"""
    
    # Slide 10: Database
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "💾 Database"
    content10.text = """• PostgreSQL database
• Foydalanuvchilar ma'lumotlari
• Guruh sozlamalari
• Video ma'lumotlari
• Mavsumlar
• Xavfsizlik ma'lumotlari"""
    
    # Slide 11: Video Yuborish Jarayoni
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "📤 Video Yuborish Jarayoni"
    content11.text = """1. Vaqt kelganda
2. Keyingi ko'rilmagan video topiladi
3. Foydalanuvchiga yuboriladi
4. Progress yangilanadi
5. Keyingi video rejalashtiriladi
6. Xatoliklar loglanadi"""
    
    # Slide 12: Foydalanuvchi Ro'yxatdan O'tkazish
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "📝 Foydalanuvchi Ro'yxatdan O'tkazish"
    content12.text = """1. Foydalanuvchi /start buyrug'ini yuboradi
2. Ma'lumotlar kiritiladi
3. Admin tasdiqlaydi
4. Ruxsat beriladi
5. Bot funksiyalari ochiladi
6. Video yuborish boshlanadi"""
    
    # Slide 13: Guruh Video Sozlamalari
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "⚙️ Guruh Video Sozlamalari"
    content13.text = """• /set_group_video centris 2
• /set_group_video golden 1
• Mavsumdan boshlash
• Guruh uchun maxsus sozlamalar
• Vaqtni o'zgartirish
• Sozlamalarni ko'rish"""
    
    # Slide 14: Mavsum Qo'shish
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title14 = slide14.shapes.title
    content14 = slide14.placeholders[1]
    
    title14.text = "📅 Mavsum Qo'shish"
    content14.text = """• /add_season centris 3
• /add_season golden 2
• Yangi mavsum yaratish
• Video ro'yxatini kiritish
• Mavsum ma'lumotlarini saqlash
• Avtomatik rejalashtirish"""
    
    # Slide 15: Xatoliklarni Tuzatish
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    title15 = slide15.shapes.title
    content15 = slide15.placeholders[1]
    
    title15.text = "🔧 Xatoliklarni Tuzatish"
    content15.text = """• Parse mode xatoliklari
• Database xatoliklari
• Scheduler xatoliklari
• Video yuborish xatoliklari
• Xavfsizlik xatoliklari
• Log fayllarini ko'rish"""
    
    # Slide 16: Statistika va Monitoring
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    title16 = slide16.shapes.title
    content16 = slide16.placeholders[1]
    
    title16.text = "📊 Statistika va Monitoring"
    content16.text = """• Foydalanuvchilar soni
• Yuborilgan videolar
• Xatoliklar soni
• Guruhlar soni
• Video progress
• Tizim holati"""
    
    # Slide 17: Foydalanish Qo'llanmasi
    slide17 = prs.slides.add_slide(prs.slide_layouts[1])
    title17 = slide17.shapes.title
    content17 = slide17.placeholders[1]
    
    title17.text = "📖 Foydalanish Qo'llanmasi"
    content17.text = """• Botni ishga tushirish
• Foydalanuvchi ro'yxatdan o'tkazish
• Video sozlamalari
• Guruh boshqaruvi
• Admin funksiyalari
• Xavfsizlik sozlamalari"""
    
    # Slide 18: Afzalliklar
    slide18 = prs.slides.add_slide(prs.slide_layouts[1])
    title18 = slide18.shapes.title
    content18 = slide18.placeholders[1]
    
    title18.text = "✨ Afzalliklar"
    content18.text = """• Avtomatik video yuborish
• Oson boshqaruv
• Xavfsizlik tizimi
• Ko'p guruh qo'llab-quvvatlash
• Real-time monitoring
• Xatoliklarni avtomatik tuzatish"""
    
    # Slide 19: Texnik Talablar
    slide19 = prs.slides.add_slide(prs.slide_layouts[1])
    title19 = slide19.shapes.title
    content19 = slide19.placeholders[1]
    
    title19.text = "⚙️ Texnik Talablar"
    content19.text = """• Python 3.8+
• PostgreSQL database
• Telegram Bot API
• APScheduler
• aiogram 2.x
• Linux server"""
    
    # Slide 20: Yakuniy
    slide20 = prs.slides.add_slide(prs.slide_layouts[1])
    title20 = slide20.shapes.title
    content20 = slide20.placeholders[1]
    
    title20.text = "🎯 Yakuniy"
    content20.text = """• Bot to'liq ishlayapti
• Barcha funksiyalar tushunarli
• Xatoliklar tuzatildi
• Prezentatsiya tayyor
• Foydalanish oson
• Texnik qo'llab-quvvatlash mavjud"""
    
    # Prezentatsiyani saqlash
    filename = "prezentatsiya_v6_sodda_funksiyalar.pptx"
    prs.save(filename)
    
    print(f"✅ Prezentatsiya yaratildi: {filename}")
    print(f"📊 Jami slaydlar: {len(prs.slides)}")
    print(f"🎯 Mavzu: Botning barcha funksiyalarini sodda va tushunarli holda")
    
    return filename

if __name__ == "__main__":
    try:
        filename = create_simple_functions_presentation()
        print(f"\n🎉 Prezentatsiya muvaffaqiyatli yaratildi!")
        print(f"📁 Fayl: {filename}")
        print(f"🚀 Endi siz botning barcha funksiyalarini sodda va tushunarli holda taqdimot qilishingiz mumkin!")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
