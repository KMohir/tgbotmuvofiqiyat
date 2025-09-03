#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot PowerPoint Presentation Creator V2
Professional diagrammalar va to'liq buyruqlar bilan
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_title_slide(prs, title, subtitle):
    """Title slide yaratish"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_content_slide(prs, title, content):
    """Content slide yaratish"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    content_shape.text = content
    
    # Content styling
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_architecture_diagram(prs):
    """Arxitektura diagrammasi"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🏗️ Bot Arxitekturasi"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Telegram API Box
    telegram_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(2.5), Inches(1.5))
    telegram_box.fill.solid()
    telegram_box.fill.fore_color.rgb = RGBColor(0, 150, 255)
    telegram_text = telegram_box.text_frame
    telegram_text.text = "Telegram\nAPI"
    telegram_text.paragraphs[0].font.size = Pt(16)
    telegram_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    telegram_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bot Box
    bot_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2), Inches(3), Inches(1.5))
    bot_box.fill.solid()
    bot_box.fill.fore_color.rgb = RGBColor(0, 102, 204)
    bot_text = bot_box.text_frame
    bot_text.text = "Centris Bot\nCore"
    bot_text.paragraphs[0].font.size = Pt(16)
    bot_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    bot_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Database Box
    db_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(2.5), Inches(1.5))
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = RGBColor(0, 150, 0)
    db_text = db_box.text_frame
    db_text.text = "PostgreSQL\nDatabase"
    db_text.paragraphs[0].font.size = Pt(16)
    db_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    db_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrows
    # Telegram to Bot
    arrow1 = slide.shapes.add_connector(1, Inches(3), Inches(2.75), Inches(3.5), Inches(2.75))
    arrow1.line.color.rgb = RGBColor(0, 0, 0)
    
    # Bot to Database
    arrow2 = slide.shapes.add_connector(1, Inches(6.5), Inches(2.75), Inches(7), Inches(2.75))
    arrow2.line.color.rgb = RGBColor(0, 0, 0)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Bot Telegram API orqali xabarlarni qabul qiladi, so'ng ma'lumotlar bazasiga murojaat qiladi"
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_database_diagram(prs):
    """Database diagrammasi"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🗄️ Ma'lumotlar Bazasi Tuzilishi"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Users table
    users_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(2.5), Inches(1.8))
    users_box.fill.solid()
    users_box.fill.fore_color.rgb = RGBColor(255, 200, 200)
    users_text = users_box.text_frame
    users_text.text = "users\n- id (PK)\n- username\n- subscription"
    users_text.paragraphs[0].font.size = Pt(14)
    users_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Groups table
    groups_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2), Inches(2.5), Inches(1.8))
    groups_box.fill.solid()
    groups_box.fill.fore_color.rgb = RGBColor(200, 255, 200)
    groups_text = groups_box.text_frame
    groups_text.text = "group_video_settings\n- group_id (PK)\n- centris_enabled\n- golden_enabled\n- send_times"
    groups_text.paragraphs[0].font.size = Pt(14)
    groups_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Videos table
    videos_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2), Inches(2.5), Inches(1.8))
    videos_box.fill.solid()
    videos_box.fill.fore_color.rgb = RGBColor(200, 200, 255)
    videos_text = videos_box.text_frame
    videos_text.text = "videos\n- id (PK)\n- season_id (FK)\n- url\n- title"
    videos_text.paragraphs[0].font.size = Pt(14)
    videos_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Relationships
    # Users to Groups
    rel1 = slide.shapes.add_connector(1, Inches(3), Inches(2.9), Inches(3.5), Inches(2.9))
    rel1.line.color.rgb = RGBColor(255, 0, 0)
    
    # Groups to Videos
    rel2 = slide.shapes.add_connector(1, Inches(6), Inches(2.9), Inches(6.5), Inches(2.9))
    rel2.line.color.rgb = RGBColor(0, 0, 255)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Asosiy jadvallar: users (foydalanuvchilar), group_video_settings (guruh sozlamalari), videos (videolar)"
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_video_flow_diagram(prs):
    """Video flow diagrammasi"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🎬 Video Tarqatish Jarayoni"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Process boxes
    # Step 1: Time Check
    step1_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(2), Inches(1))
    step1_box.fill.solid()
    step1_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
    step1_text = step1_box.text_frame
    step1_text.text = "1. Vaqt\nTekshirish"
    step1_text.paragraphs[0].font.size = Pt(14)
    step1_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    step1_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Step 2: Get Groups
    step2_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2), Inches(2), Inches(1))
    step2_box.fill.solid()
    step2_box.fill.fore_color.rgb = RGBColor(0, 255, 255)
    step2_text = step2_box.text_frame
    step2_text.text = "2. Guruhlarni\nOlish"
    step2_text.paragraphs[0].font.size = Pt(14)
    step2_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    step2_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Step 3: Send Videos
    step3_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2), Inches(2), Inches(1))
    step3_box.fill.solid()
    step3_box.fill.fore_color.rgb = RGBColor(255, 0, 255)
    step3_text = step3_box.text_frame
    step3_text.text = "3. Videolarni\nYuborish"
    step3_text.paragraphs[0].font.size = Pt(14)
    step3_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    step3_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Step 4: Update Progress
    step4_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(2), Inches(2), Inches(1))
    step4_box.fill.solid()
    step4_box.fill.fore_color.rgb = RGBColor(0, 255, 0)
    step4_text = step4_box.text_frame
    step4_text.text = "4. Progress\nYangilash"
    step4_text.paragraphs[0].font.size = Pt(14)
    step4_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    step4_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrows
    arrow1 = slide.shapes.add_connector(1, Inches(2.5), Inches(2.5), Inches(3), Inches(2.5))
    arrow1.line.color.rgb = RGBColor(0, 0, 0)
    
    arrow2 = slide.shapes.add_connector(1, Inches(5), Inches(2.5), Inches(5.5), Inches(2.5))
    arrow2.line.color.rgb = RGBColor(0, 0, 0)
    
    arrow3 = slide.shapes.add_connector(1, Inches(7.5), Inches(2.5), Inches(8), Inches(2.5))
    arrow3.line.color.rgb = RGBColor(0, 0, 0)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Video yuborish 4 ta asosiy qadamdan iborat: vaqt tekshirish → guruhlarni olish → videolarni yuborish → progress yangilash"
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_commands_slide(prs):
    """Bot buyruqlari haqida batafsil ma'lumot"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📱 Bot Buyruqlari - Batafsil Ma'lumot"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Commands in boxes
    commands = [
        ("/start", "Botni ishga tushirish, foydalanuvchi ro'yxatdan o'tkazish", RGBColor(255, 200, 200)),
        ("/centris_towers", "Centris Towers loyihasi videolarini ko'rish", RGBColor(200, 255, 200)),
        ("/golden_lake", "Golden Lake loyihasi videolarini ko'rish", RGBColor(200, 200, 255)),
        ("/set_group_video", "Guruh uchun video sozlamalari (admin)", RGBColor(255, 255, 200)),
        ("/show_group_video_settings", "Guruh sozlamalarini ko'rish (admin)", RGBColor(255, 200, 255)),
        ("/update_video_progress", "Video progress yangilash (admin)", RGBColor(200, 255, 255)),
        ("/send_specific_video", "Maxsus video yuborish (admin)", RGBColor(255, 255, 255))
    ]
    
    y_position = 2
    for i, (cmd, desc, color) in enumerate(commands):
        # Command box
        cmd_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_position), Inches(2.5), Inches(0.6))
        cmd_box.fill.solid()
        cmd_box.fill.fore_color.rgb = color
        cmd_text = cmd_box.text_frame
        cmd_text.text = cmd
        cmd_text.paragraphs[0].font.size = Pt(16)
        cmd_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cmd_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description box
        desc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(y_position), Inches(6.3), Inches(0.6))
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        desc_text = desc_box.text_frame
        desc_text.text = desc
        desc_text.paragraphs[0].font.size = Pt(14)
        desc_text.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        y_position += 0.8
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Barcha buyruqlar o'zbek tilida va tushunarli. Admin buyruqlari faqat ruxsat berilgan foydalanuvchilar uchun."
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Asosiy funksiya"""
    print("🎬 Centris Bot PowerPoint Presentation V2 yaratilmoqda...")
    
    # Yangi presentation yaratish
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "🎬 CENTRIS TOWERS & GOLDEN LAKE",
        "Telegram Bot Arxitekturasi va Ishlash Prinsipi\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "📋 Loyiha Haqida",
        """• Centris Towers va Golden Lake loyihalari uchun Telegram bot
• Avtomatik video tarqatish tizimi
• Admin panel va foydalanuvchi interfeysi
• PostgreSQL ma'lumotlar bazasi
• APScheduler bilan vaqt rejalashtirish
• Xavfsizlik tizimi va whitelist
• Professional arxitektura va kod sifat"""
    )
    
    # Slide 3: Architecture Diagram
    create_architecture_diagram(prs)
    
    # Slide 4: Database Schema
    create_database_diagram(prs)
    
    # Slide 5: Video Distribution Flow
    create_video_flow_diagram(prs)
    
    # Slide 6: Commands Overview
    create_commands_slide(prs)
    
    # Slide 7: Security System
    create_content_slide(
        prs,
        "🔒 Xavfsizlik Tizimi",
        """• Whitelist tizimi - faqat ruxsat berilgan guruhlar
• Admin roli - faqat admin buyruqlarini bajaradi
• Super-admin roli - barcha huquqlarga ega
• Group verification - guruh mavjudligini tekshirish
• Rate limiting - haddan tashqari so'rovlarni cheklash
• Input validation - kiritilgan ma'lumotlarni tekshirish
• Xavfsizlik choralari to'liq amalga oshirilgan"""
    )
    
    # Slide 8: Technical Stack
    create_content_slide(
        prs,
        "🛠️ Texnik Stack",
        """• Python 3.13 - asosiy dasturlash tili
• aiogram 2.x - Telegram Bot API kutubxonasi
• PostgreSQL - ma'lumotlar bazasi
• APScheduler - vaqt rejalashtirish
• psycopg2 - PostgreSQL connector
• pytz - vaqt zonasi boshqaruvi
• aiohttp - HTTP client kutubxonasi
• Professional kod tuzilishi va xatoliklarni bartaraf etish"""
    )
    
    # Slide 9: Video Scheduling System
    create_content_slide(
        prs,
        "⏰ Video Rejalashtirish Tizimi",
        """• APScheduler bilan professional vaqt boshqaruvi
• Har bir guruh uchun alohida vaqt sozlamalari
• Centris Towers: 08:00 va 20:00
• Golden Lake: 11:00 va 23:00
• Avtomatik progress kuzatish
• Xatoliklar va muammolarni bartaraf etish
• To'liq ishonchli va barqaror tizim"""
    )
    
    # Slide 10: Admin Panel Features
    create_content_slide(
        prs,
        "⚙️ Admin Panel Xususiyatlari",
        """• Guruh sozlamalarini boshqarish
• Video progress yangilash
• Maxsus video yuborish
• Guruh nomlarini yangilash
• Barcha rejalashtirilgan videolarni yuborish
• Test video yuborish
• Professional admin interfeysi
• Xavfsizlik va nazorat"""
    )
    
    # Slide 11: User Experience
    create_content_slide(
        prs,
        "👥 Foydalanuvchi Tajribasi",
        """• Oson va tushunarli interfeys
• O'zbek tilida to'liq qo'llab-quvvatlash
• Avtomatik video tarqatish
• Progress kuzatish
• Qulay buyruqlar
• Tezkor javob berish
• Professional xizmat sifat"""
    )
    
    # Slide 12: Future Plans
    create_content_slide(
        prs,
        "🚀 Kelajak Rejalari",
        """• Yangi loyihalar qo'shish
• Video analytics va statistika
• Avtomatik content generation
• AI-powered video recommendations
• Mobile app development
• Integration with CRM systems
• Multi-language support
• Advanced admin dashboard
• Cloud deployment va scaling"""
    )
    
    # Slide 13: Conclusion
    create_content_slide(
        prs,
        "✅ Xulosa va Natijalar",
        """• Centris Bot muvaffaqiyatli ishlayapti
• Barcha xatolar bartaraf etildi
• Video tarqatish tizimi to'liq ishlayapti
• Admin panel qulay va funksional
• Xavfsizlik tizimi ishonchli
• Kod sifatli va tuzilgan
• Professional arxitektura
• Kelajakda rivojlantirish imkoniyatlari keng
• Mijozlar mamnun"""
    )
    
    # Slide 14: Thank You
    create_title_slide(
        prs,
        "🎉 Rahmat!",
        "Centris Bot haqida so'rashganingiz uchun!\n\nSavollar bo'lsa, so'rang!\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Faylni saqlash
    filename = "prezentatsiya_v2.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint presentation V2 muvaffaqiyatli yaratildi: {filename}")
    print(f"📊 Jami slaydlar soni: {len(prs.slides)}")
    print(f"🎯 Professional diagrammalar qo'shildi")
    print(f"📱 Barcha buyruqlar haqida batafsil ma'lumot")
    print(f"🇺🇿 O'zbek tilida tayyorlandi")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎬 {filename} fayli muvaffaqiyatli yaratildi!")
        print("📁 Fayl joriy papkada joylashgan")
        print("🎯 Professional diagrammalar va to'liq ma'lumotlar bilan")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
        print("💡 python-pptx kutubxonasi o'rnatilganligini tekshiring")
