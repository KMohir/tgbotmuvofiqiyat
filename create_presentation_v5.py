#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot PowerPoint Presentation Creator V5
/set_group_video va /add_season buyruqlari haqida to'liq tushuntirish
prezentatsiya_v2.pptx ga o'xshash ko'rinishda
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    """Title slide yaratish"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_content_slide(prs, title, content):
    """Content slide yaratish"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    content_shape.text = content
    
    # Content styling
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_command_detailed_slide(prs, command, description, functions, examples, workflow):
    """Buyruq haqida batafsil ma'lumot slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"📱 {command} - Batafsil Ma'lumot"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(18)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Functions
    func_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(3.8), Inches(2.5))
    func_frame = func_box.text_frame
    func_frame.text = f"🔹 **Funksiyalar:**\n{functions}"
    func_frame.paragraphs[0].font.size = Pt(16)
    func_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    func_frame.paragraphs[0].font.bold = True
    
    for paragraph in func_frame.paragraphs[1:]:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    # Examples
    ex_box = slide.shapes.add_textbox(Inches(5.2), Inches(3.2), Inches(3.8), Inches(2.5))
    ex_frame = ex_box.text_frame
    ex_frame.text = f"💡 **Misollar:**\n{examples}"
    ex_frame.paragraphs[0].font.size = Pt(16)
    ex_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    ex_frame.paragraphs[0].font.bold = True
    
    for paragraph in ex_frame.paragraphs[1:]:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    # Workflow
    workflow_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1.5))
    workflow_frame = workflow_box.text_frame
    workflow_frame.text = f"🔄 **Ishlash Jarayoni:**\n{workflow}"
    workflow_frame.paragraphs[0].font.size = Pt(16)
    workflow_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    workflow_frame.paragraphs[0].font.bold = True
    
    for paragraph in workflow_frame.paragraphs[1:]:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_workflow_diagram_slide(prs, title, workflow_steps):
    """Workflow diagramma slide yaratish (text bilan)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Workflow steps
    workflow_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    workflow_frame = workflow_box.text_frame
    workflow_frame.text = workflow_steps
    
    # Content styling
    for paragraph in workflow_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("1️⃣") or paragraph.text.startswith("2️⃣") or paragraph.text.startswith("3️⃣") or paragraph.text.startswith("4️⃣") or paragraph.text.startswith("5️⃣"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_admin_commands_overview_slide(prs):
    """Admin buyruqlari haqida umumiy ma'lumot"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "⚙️ Admin Buyruqlari - Umumiy Ko'rinish"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Admin commands overview
    commands_text = """🔹 **Guruh Boshqaruvi:**
   • /set_group_video - Video sozlamalari belgilash
   • /show_group_video_settings - Sozlamalarni ko'rish
   • /update_group_names - Guruh nomlarini yangilash

🔹 **Video Boshqaruvi:**
   • /add_season - Yangi mavsum qo'shish
   • /send_specific_video - Maxsus video yuborish
   • /update_video_progress - Progress yangilash

🔹 **Sistem Boshqaruvi:**
   • /send_all_planned_videos - Barcha videolarni yuborish
   • /test_send_video_all_groups - Test video yuborish
   • /admin_help - Admin buyruqlari haqida yordam

🔹 **Xavfsizlik:**
   • Admin huquqini tekshirish
   • Guruh ruxsatini tekshirish
   • Foydalanuvchi huquqlarini boshqarish"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = commands_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_security_system_slide(prs):
    """Xavfsizlik tizimi slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🔒 Xavfsizlik Tizimi va Admin Huquqlari"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Security system
    security_text = """🔐 **Admin Huquqlari:**
   • **is_admin()** - Oddiy admin huquqini tekshirish
   • **is_super_admin()** - Super-admin huquqini tekshirish
   • **check_group_whitelist()** - Guruh ruxsatini tekshirish

🛡️ **Xavfsizlik Choralari:**
   • Faqat ruxsat berilgan guruhlar
   • Admin buyruqlarini himoya qilish
   • Input validation va sanitization
   • Rate limiting va spam himoyasi

🔑 **Huquq Darajalari:**
   • **Foydalanuvchi** - Faqat video ko'rish
   • **Admin** - Guruh sozlamalari, video yuborish
   • **Super-Admin** - Mavsum qo'shish, barcha sozlamalar

💻 **Xavfsizlik Funksiyalari:**
   • check_user_permissions()
   • validate_admin_command()
   • secure_database_query()
   • audit_log_admin_actions()"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = security_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔐") or paragraph.text.startswith("🛡️") or paragraph.text.startswith("🔑") or paragraph.text.startswith("💻"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_database_integration_slide(prs):
    """Database integratsiyasi slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🗄️ Database Integratsiyasi va So'rovlar"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Database integration
    db_text = """📊 **Asosiy Jadval:**
   • **users** - Foydalanuvchilar ma'lumotlari
   • **group_video_settings** - Guruh video sozlamalari
   • **videos** - Video ma'lumotlari va URL lar
   • **seasons** - Mavsum ma'lumotlari

🔍 **SET_GROUP_VIDEO So'rovlari:**
   • UPDATE group_video_settings SET centris_season = %s
   • UPDATE group_video_settings SET golden_season = %s
   • INSERT INTO group_video_settings (group_id, group_name, ...)

➕ **ADD_SEASON So'rovlari:**
   • INSERT INTO videos (season_id, url, title, project_type)
   • INSERT INTO seasons (project_type, season_number)
   • UPDATE video_progress SET total_videos = total_videos + 1

🔗 **Bog'lanishlar:**
   • group_video_settings ↔ videos (One-to-Many)
   • videos ↔ seasons (Many-to-One)
   • users ↔ group_video_settings (Many-to-Many)

💾 **PostgreSQL Xususiyatlari:**
   • psycopg2 connector
   • Connection pooling
   • Transaction management
   • Error handling va rollback"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = db_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("📊") or paragraph.text.startswith("🔍") or paragraph.text.startswith("➕") or paragraph.text.startswith("🔗") or paragraph.text.startswith("💾"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def main():
    """Asosiy funksiya"""
    print("🎬 Centris Bot PowerPoint Presentation V5 yaratilmoqda...")
    print("📱 /set_group_video va /add_season buyruqlari haqida to'liq tushuntirish")
    print("🎯 prezentatsiya_v2.pptx ga o'xshash ko'rinishda")
    
    # Yangi presentation yaratish
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "🎬 CENTRIS TOWERS & GOLDEN LAKE",
        "Admin Buyruqlari va Funksiyalari\n\n/set_group_video va /add_season haqida to'liq ma'lumot\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "📋 Loyiha va Admin Buyruqlari Haqida",
        """• Centris Towers va Golden Lake loyihalari uchun Telegram bot
• Admin panel orqali guruh va video boshqaruvi
• /set_group_video - Guruh uchun video sozlamalari belgilash
• /add_season - Yangi mavsum va videolar qo'shish
• Xavfsizlik tizimi va admin huquqlari
• PostgreSQL ma'lumotlar bazasi integratsiyasi
• Professional arxitektura va kod sifat
• Barcha admin funksiyalari to'liq ko'rsatildi"""
    )
    
    # Slide 3: Admin Commands Overview
    create_admin_commands_overview_slide(prs)
    
    # Slide 4: SET_GROUP_VIDEO Command Detailed
    create_command_detailed_slide(
        prs,
        "/set_group_video",
        "Guruh uchun video sozlamalarini belgilash va boshqarish",
        """• set_group_video_handler() - Asosiy handler
• update_group_settings() - Sozlamalarni yangilash
• check_admin_permissions() - Admin huquqini tekshirish
• validate_group_exists() - Guruh mavjudligini tekshirish
• update_database_settings() - Database ni yangilash
• send_confirmation_message() - Tasdiqlash xabari""",
        """• /set_group_video centris 2
• /set_group_video golden 3
• /set_group_video centris 1 golden 4
• /set_group_video all 5""",
        """1️⃣ Admin huquqini tekshirish
2️⃣ Guruh mavjudligini tekshirish
3️⃣ Loyiha va mavsum parametrlarini olish
4️⃣ Database sozlamalarini yangilash
5️⃣ Tasdiqlash xabarini yuborish"""
    )
    
    # Slide 5: ADD_SEASON Command Detailed
    create_command_detailed_slide(
        prs,
        "/add_season",
        "Yangi mavsum va videolar qo'shish (Super-Admin)",
        """• add_season_handler() - Asosiy handler
• check_super_admin_permissions() - Super-admin huquqini tekshirish
• validate_video_url() - Video URL ni tekshirish
• insert_new_season() - Yangi mavsum qo'shish
• insert_video_records() - Video yozuvlarini qo'shish
• update_progress_tracking() - Progress kuzatishni yangilash""",
        """• /add_season centris 5 https://example.com/video.mp4
• /add_season golden 4 https://example.com/golden.mp4
• /add_season centris 6 https://example.com/centris6.mp4
• /add_season golden 5 https://example.com/golden5.mp4""",
        """1️⃣ Super-admin huquqini tekshirish
2️⃣ Video URL ni tekshirish va validatsiya qilish
3️⃣ Yangi mavsum yozuvini qo'shish
4️⃣ Video ma'lumotlarini database ga saqlash
5️⃣ Progress kuzatishni yangilash va tasdiqlash"""
    )
    
    # Slide 6: SET_GROUP_VIDEO Workflow
    create_workflow_diagram_slide(
        prs,
        "🔄 /set_group_video - Ishlash Jarayoni",
        """1️⃣ **Admin Buyruqni Yuboradi:**
   /set_group_video centris 2

2️⃣ **Huquqni Tekshirish:**
   is_admin() → True/False
   check_group_whitelist() → Guruh ruxsati

3️⃣ **Parametrlarni Olish:**
   project = "centris"
   season = 2
   group_id = current_group_id

4️⃣ **Database Yangilash:**
   UPDATE group_video_settings 
   SET centris_season = 2 
   WHERE group_id = %s

5️⃣ **Tasdiqlash Xabari:**
   "✅ Centris Towers mavsum 2 dan boshlab videolar yuboriladi" """
    )
    
    # Slide 7: ADD_SEASON Workflow
    create_workflow_diagram_slide(
        prs,
        "🔄 /add_season - Ishlash Jarayoni",
        """1️⃣ **Super-Admin Buyruqni Yuboradi:**
   /add_season centris 5 https://example.com/video.mp4

2️⃣ **Super-Admin Huquqini Tekshirish:**
   is_super_admin() → True/False
   check_super_admin_permissions() → Huquq mavjudligi

3️⃣ **Video URL ni Tekshirish:**
   validate_video_url() → URL to'g'ri/noto'g'ri
   check_video_accessibility() → Video mavjudligi

4️⃣ **Database ga Qo'shish:**
   INSERT INTO seasons (project_type, season_number)
   INSERT INTO videos (season_id, url, title, project_type)

5️⃣ **Progress Yangilash:**
   UPDATE video_progress SET total_videos = total_videos + 1
   "✅ Centris Towers mavsum 5 qo'shildi, 1 ta video" """
    )
    
    # Slide 8: Security System
    create_security_system_slide(prs)
    
    # Slide 9: Database Integration
    create_database_integration_slide(prs)
    
    # Slide 10: Command Parameters and Validation
    create_content_slide(
        prs,
        "🔍 Buyruq Parametrlari va Validatsiya",
        """🔹 **/set_group_video parametrlari:**
   • project: "centris" yoki "golden" yoki "all"
   • season_number: 1, 2, 3, 4, 5... (integer)
   • group_id: Avtomatik olinadi (current group)

🔹 **/add_season parametrlari:**
   • project: "centris" yoki "golden"
   • season_number: Yangi mavsum raqami
   • video_url: Video fayl URL si (https://...)

🔹 **Validatsiya Qoidalari:**
   • Admin huquqini tekshirish
   • Guruh mavjudligini tekshirish
   • Video URL to'g'riligini tekshirish
   • Mavsum raqami mantiqiy ekanligini tekshirish
   • Database xatoliklarini qayta ishlash"""
    )
    
    # Slide 11: Error Handling and User Feedback
    create_content_slide(
        prs,
        "⚠️ Xatoliklarni Qayta Ishlash va Foydalanuvchi Javobi",
        """🔹 **Xatolik Turlari:**
   • "Siz admin emassiz!" - Huquq yo'q
   • "Guruh topilmadi!" - Guruh mavjud emas
   • "Noto'g'ri URL!" - Video URL noto'g'ri
   • "Mavsum allaqachon mavjud!" - Duplicate mavsum
   • "Database xatoligi!" - Database muammosi

🔹 **Muvaffaqiyat Xabarlari:**
   • "✅ Centris Towers mavsum 2 dan boshlab videolar yuboriladi"
   • "✅ Yangi mavsum 5 qo'shildi, 1 ta video"
   • "✅ Sozlamalar muvaffaqiyatli yangilandi"

🔹 **Xatolik Qayta Ishlash:**
   • try-catch bloklari
   • User-friendly xatolik xabarlari
   • Logging va monitoring
   • Automatic retry mechanisms"""
    )
    
    # Slide 12: Integration with Video System
    create_content_slide(
        prs,
        "🎬 Video Tizimi bilan Integratsiya",
        """🔹 **Video Yuborish Jarayoni:**
   • /set_group_video → Video sozlamalari yangilanadi
   • Scheduler → Yangi sozlamalar bo'yicha videolar yuboriladi
   • APScheduler → Vaqt bo'yicha avtomatik yuborish
   • Progress tracking → Video yuborish holati kuzatiladi

🔹 **Mavsum Qo'shish Natijasi:**
   • /add_season → Yangi mavsum va videolar qo'shiladi
   • Video klaviaturasi → Yangi mavsum ko'rsatiladi
   • Foydalanuvchilar → Yangi videolarni ko'ra oladi
   • Admin panel → Yangi mavsum boshqariladi

🔹 **Avtomatik Jarayonlar:**
   • Video rejalashtirish
   • Progress kuzatish
   • Error handling
   • User notifications"""
    )
    
    # Slide 13: Future Enhancements
    create_content_slide(
        prs,
        "🚀 Kelajak Yaxshilanishlar va Yangi Funksiyalar",
        """🔹 **Admin Panel Yaxshilanishlari:**
   • Bulk operations - Bir necha guruhni bir vaqtda
   • Advanced filtering - Guruh va mavsum bo'yicha filtrlash
   • Statistics dashboard - Video yuborish statistikasi
   • User management - Foydalanuvchilarni boshqarish

🔹 **Video Boshqaruvi:**
   • Video categories - Video kategoriyalari
   • Content scheduling - Content rejalashtirish
   • Quality control - Video sifatini nazorat qilish
   • Backup systems - Zavol tizimlari

🔹 **Xavfsizlik:**
   • Two-factor authentication - Ikki bosqichli autentifikatsiya
   • Audit logging - Barcha harakatlarni qayd qilish
   • Role-based access control - Rol asosida huquq boshqaruvi
   • API rate limiting - API so'rovlarini cheklash"""
    )
    
    # Slide 14: Conclusion
    create_content_slide(
        prs,
        "✅ Xulosa va Natijalar",
        """• /set_group_video buyrug'i to'liq ishlayapti
• /add_season buyrug'i super-admin uchun tayyor
• Barcha admin funksiyalari xavfsiz va ishonchli
• Database integratsiyasi professional
• Xatoliklarni qayta ishlash tizimi to'liq
• Foydalanuvchi interfeysi qulay va tushunarli
• Video tizimi bilan to'liq integratsiya
• Kelajakda rivojlantirish imkoniyatlari keng
• Admin panel professional va funksional"""
    )
    
    # Slide 15: Thank You
    create_title_slide(
        prs,
        "🎉 Rahmat!",
        "Admin buyruqlari haqida so'rashganingiz uchun!\n\n/set_group_video va /add_season to'liq tushuntirildi!\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Faylni saqlash
    filename = "prezentatsiya_v5_admin_buyruqlar.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint presentation V5 muvaffaqiyatli yaratildi: {filename}")
    print(f"📊 Jami slaydlar soni: {len(prs.slides)}")
    print(f"📱 /set_group_video va /add_season buyruqlari to'liq tushuntirildi")
    print(f"🔹 Barcha funksiyalar va parametrlar ko'rsatildi")
    print(f"🇺🇿 O'zbek tilida tayyorlandi")
    print(f"💡 prezentatsiya_v2.pptx ga o'xshash professional ko'rinish")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎬 {filename} fayli muvaffaqiyatli yaratildi!")
        print("📁 Fayl joriy papkada joylashgan")
        print("📱 Admin buyruqlari haqida to'liq ma'lumot")
        print("💡 Professional ko'rinish va tuzilish")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
        print("💡 python-pptx kutubxonasi o'rnatilganligini tekshiring")
