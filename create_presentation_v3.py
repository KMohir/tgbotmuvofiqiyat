#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot PowerPoint Presentation Creator V3
Rasm diagrammalar bilan professional ko'rinish
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    """Title slide yaratish"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_content_slide(prs, title, content):
    """Content slide yaratish"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    content_shape.text = content
    
    # Content styling
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_diagram_slide_with_text(prs, title, diagram_text, description):
    """Diagramma slide yaratish (text diagramma bilan)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Diagram text (temporary, will be replaced with image)
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    diagram_frame = diagram_box.text_frame
    diagram_frame.text = f"📊 {diagram_text}\n\n💡 Bu joyda professional diagramma rasm ko'rinishida bo'ladi"
    diagram_frame.paragraphs[0].font.size = Pt(18)
    diagram_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    diagram_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_commands_slide(prs):
    """Bot buyruqlari haqida batafsil ma'lumot"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📱 Bot Buyruqlari - Batafsil Ma'lumot"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Commands in text format
    commands_text = """🔹 /start - Botni ishga tushirish, foydalanuvchi ro'yxatdan o'tkazish

🔹 /centris_towers - Centris Towers loyihasi videolarini ko'rish

🔹 /golden_lake - Golden Lake loyihasi videolarini ko'rish

🔹 /set_group_video - Guruh uchun video sozlamalari (ADMIN)

🔹 /show_group_video_settings - Guruh sozlamalarini ko'rish (ADMIN)

🔹 /update_video_progress - Video progress yangilash (ADMIN)

🔹 /send_specific_video - Maxsus video yuborish (ADMIN)"""
    
    # Commands content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = commands_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Barcha buyruqlar o'zbek tilida va tushunarli. Admin buyruqlari faqat ruxsat berilgan foydalanuvchilar uchun."
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_image_diagram_instructions(prs):
    """Rasm diagrammalar qo'yish bo'yicha ko'rsatmalar"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🖼️ Diagrammalarni Rasm Ko'rinishida Qo'yish"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Instructions
    instructions_text = """📋 Diagrammalarni rasm ko'rinishida qo'yish uchun:

1️⃣ **draw.io** yoki **Mermaid** da diagramma yarating
2️⃣ **PNG/JPG** formatida saqlang
3️⃣ PowerPoint da **Insert → Picture** orqali qo'shing
4️⃣ Diagrammani slaydga joylashtiring

🎯 **Kerakli diagrammalar:**
• Bot arxitekturasi
• Database schema  
• Video tarqatish jarayoni
• User flow
• Admin panel workflow

💡 **Professional ko'rinish uchun:**
• Yuqori sifatli rasm (300 DPI)
• Oq fon bilan
• Aniq va tushunarli
• Rangli va chiroyli"""
    
    # Instructions content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = instructions_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("1️⃣") or paragraph.text.startswith("🎯") or paragraph.text.startswith("💡"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def main():
    """Asosiy funksiya"""
    print("🎬 Centris Bot PowerPoint Presentation V3 yaratilmoqda...")
    print("🖼️ Rasm diagrammalar bilan professional ko'rinish")
    
    # Yangi presentation yaratish
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "🎬 CENTRIS TOWERS & GOLDEN LAKE",
        "Telegram Bot Arxitekturasi va Ishlash Prinsipi\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "📋 Loyiha Haqida",
        """• Centris Towers va Golden Lake loyihalari uchun Telegram bot
• Avtomatik video tarqatish tizimi
• Admin panel va foydalanuvchi interfeysi
• PostgreSQL ma'lumotlar bazasi
• APScheduler bilan vaqt rejalashtirish
• Xavfsizlik tizimi va whitelist
• Professional arxitektura va kod sifat"""
    )
    
    # Slide 3: Architecture Diagram (placeholder for image)
    create_diagram_slide_with_text(
        prs,
        "🏗️ Bot Arxitekturasi",
        "Bot qanday ishlayotgani va qaysi komponentlar bilan bog'langan",
        "Bu joyda professional arxitektura diagrammasi rasm ko'rinishida bo'ladi"
    )
    
    # Slide 4: Database Schema (placeholder for image)
    create_diagram_slide_with_text(
        prs,
        "🗄️ Ma'lumotlar Bazasi Tuzilishi",
        "Ma'lumotlar bazasida qanday jadvallar va ular o'rtasidagi bog'lanishlar",
        "Bu joyda professional database schema diagrammasi rasm ko'rinishida bo'ladi"
    )
    
    # Slide 5: Video Distribution Flow (placeholder for image)
    create_diagram_slide_with_text(
        prs,
        "🎬 Video Tarqatish Jarayoni",
        "Video qanday vaqtda va qanday tartibda guruhlarga yuboriladi",
        "Bu joyda professional video flow diagrammasi rasm ko'rinishida bo'ladi"
    )
    
    # Slide 6: Commands Overview
    create_commands_slide(prs)
    
    # Slide 7: Security System
    create_content_slide(
        prs,
        "🔒 Xavfsizlik Tizimi",
        """• Whitelist tizimi - faqat ruxsat berilgan guruhlar
• Admin roli - faqat admin buyruqlarini bajaradi
• Super-admin roli - barcha huquqlarga ega
• Group verification - guruh mavjudligini tekshirish
• Rate limiting - haddan tashqari so'rovlarni cheklash
• Input validation - kiritilgan ma'lumotlarni tekshirish
• Xavfsizlik choralari to'liq amalga oshirilgan"""
    )
    
    # Slide 8: Technical Stack
    create_content_slide(
        prs,
        "🛠️ Texnik Stack",
        """• Python 3.13 - asosiy dasturlash tili
• aiogram 2.x - Telegram Bot API kutubxonasi
• PostgreSQL - ma'lumotlar bazasi
• APScheduler - vaqt rejalashtirish
• psycopg2 - PostgreSQL connector
• pytz - vaqt zonasi boshqaruvi
• aiohttp - HTTP client kutubxonasi
• Professional kod tuzilishi va xatoliklarni bartaraf etish"""
    )
    
    # Slide 9: Video Scheduling System
    create_content_slide(
        prs,
        "⏰ Video Rejalashtirish Tizimi",
        """• APScheduler bilan professional vaqt boshqaruvi
• Har bir guruh uchun alohida vaqt sozlamalari
• Centris Towers: 08:00 va 20:00
• Golden Lake: 11:00 va 23:00
• Avtomatik progress kuzatish
• Xatoliklar va muammolarni bartaraf etish
• To'liq ishonchli va barqaror tizim"""
    )
    
    # Slide 10: Admin Panel Features
    create_content_slide(
        prs,
        "⚙️ Admin Panel Xususiyatlari",
        """• Guruh sozlamalarini boshqarish
• Video progress yangilash
• Maxsus video yuborish
• Guruh nomlarini yangilash
• Barcha rejalashtirilgan videolarni yuborish
• Test video yuborish
• Professional admin interfeysi
• Xavfsizlik va nazorat"""
    )
    
    # Slide 11: User Experience
    create_content_slide(
        prs,
        "👥 Foydalanuvchi Tajribasi",
        """• Oson va tushunarli interfeys
• O'zbek tilida to'liq qo'llab-quvvatlash
• Avtomatik video tarqatish
• Progress kuzatish
• Qulay buyruqlar
• Tezkor javob berish
• Professional xizmat sifat"""
    )
    
    # Slide 12: Future Plans
    create_content_slide(
        prs,
        "🚀 Kelajak Rejalari",
        """• Yangi loyihalar qo'shish
• Video analytics va statistika
• Avtomatik content generation
• AI-powered video recommendations
• Mobile app development
• Integration with CRM systems
• Multi-language support
• Advanced admin dashboard
• Cloud deployment va scaling"""
    )
    
    # Slide 13: Image Diagram Instructions
    create_image_diagram_instructions(prs)
    
    # Slide 14: Conclusion
    create_content_slide(
        prs,
        "✅ Xulosa va Natijalar",
        """• Centris Bot muvaffaqiyatli ishlayapti
• Barcha xatolar bartaraf etildi
• Video tarqatish tizimi to'liq ishlayapti
• Admin panel qulay va funksional
• Xavfsizlik tizimi ishonchli
• Kod sifatli va tuzilgan
• Professional arxitektura
• Kelajakda rivojlantirish imkoniyatlari keng
• Mijozlar mamnun"""
    )
    
    # Slide 15: Thank You
    create_title_slide(
        prs,
        "🎉 Rahmat!",
        "Centris Bot haqida so'rashganingiz uchun!\n\nSavollar bo'lsa, so'rang!\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Faylni saqlash
    filename = "prezentatsiya_v3.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint presentation V3 muvaffaqiyatli yaratildi: {filename}")
    print(f"📊 Jami slaydlar soni: {len(prs.slides)}")
    print(f"🖼️ Rasm diagrammalar uchun joylar tayyorlandi")
    print(f"📱 Barcha buyruqlar haqida batafsil ma'lumot")
    print(f"🇺🇿 O'zbek tilida tayyorlandi")
    print(f"💡 Diagrammalarni rasm ko'rinishida qo'yish uchun ko'rsatmalar qo'shildi")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎬 {filename} fayli muvaffaqiyatli yaratildi!")
        print("📁 Fayl joriy papkada joylashgan")
        print("🖼️ Rasm diagrammalar uchun joylar tayyor")
        print("💡 Ko'rsatmalar bilan professional ko'rinish")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
        print("💡 python-pptx kutubxonasi o'rnatilganligini tekshiring")
