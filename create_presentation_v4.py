#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot PowerPoint Presentation Creator V4
Diagrammalarsiz, funksiyalar va struktura bilan
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    """Title slide yaratish"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_content_slide(prs, title, content):
    """Content slide yaratish"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    content_shape.text = content
    
    # Content styling
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_functions_slide(prs, title, functions_list):
    """Funksiyalar ro'yxati slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Functions content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = functions_list
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_code_structure_slide(prs, title, structure_info):
    """Kod struktura slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Structure content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = structure_info
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("📁") or paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_database_structure_slide(prs):
    """Database struktura slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🗄️ Database Struktura va Jadval Tuzilishi"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Database structure
    db_structure = """📊 **users jadvali:**
🔹 id (SERIAL PRIMARY KEY) - Foydalanuvchi ID
🔹 username (VARCHAR) - Telegram username
🔹 subscription (BOOLEAN) - Obuna holati
🔹 created_at (TIMESTAMP) - Yaratilgan vaqt

📊 **group_video_settings jadvali:**
🔹 group_id (BIGINT PRIMARY KEY) - Guruh ID
🔹 group_name (VARCHAR) - Guruh nomi
🔹 centris_enabled (BOOLEAN) - Centris loyihasi yoqilgan
🔹 golden_enabled (BOOLEAN) - Golden loyihasi yoqilgan
🔹 centris_season (INTEGER) - Centris mavsumi
🔹 golden_season (INTEGER) - Golden mavsumi
🔹 send_times (TEXT) - Yuborish vaqtlari

📊 **videos jadvali:**
🔹 id (SERIAL PRIMARY KEY) - Video ID
🔹 season_id (INTEGER) - Mavsum ID
🔹 url (TEXT) - Video URL
🔹 title (VARCHAR) - Video nomi
🔹 project_type (VARCHAR) - Loyiha turi

📊 **Bog'lanishlar:**
🔹 users ↔ group_video_settings (Many-to-Many)
🔹 group_video_settings ↔ videos (One-to-Many)"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = db_structure
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("📊") or paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_bot_commands_detailed_slide(prs):
    """Bot buyruqlari haqida batafsil ma'lumot"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📱 Bot Buyruqlari - Barcha Funksiyalar"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Commands in text format
    commands_text = """🔹 **/start** - Botni ishga tushirish, foydalanuvchi ro'yxatdan o'tkazish
   Funksiya: register_user(), check_user_exists()

🔹 **/centris_towers** - Centris Towers loyihasi videolarini ko'rish
   Funksiya: get_video_keyboard_from_db(), send_video()

🔹 **/golden_lake** - Golden Lake loyihasi videolarini ko'rish
   Funksiya: get_video_keyboard_from_db(), send_video()

🔹 **/set_group_video** - Guruh uchun video sozlamalari (ADMIN)
   Funksiya: set_group_video_handler(), update_group_settings()

🔹 **/show_group_video_settings** - Guruh sozlamalarini ko'rish (ADMIN)
   Funksiya: show_group_video_settings(), get_all_groups_with_settings()

🔹 **/update_video_progress** - Video progress yangilash (ADMIN)
   Funksiya: update_video_progress(), update_season_progress()

🔹 **/send_specific_video** - Maxsus video yuborish (ADMIN)
   Funksiya: send_specific_video(), send_video_to_group()

🔹 **/send_all_planned_videos** - Barcha rejalashtirilgan videolarni yuborish (ADMIN)
   Funksiya: send_all_planned_videos(), process_all_groups()

🔹 **/test_send_video_all_groups** - Test video yuborish (ADMIN)
   Funksiya: test_send_video_all_groups(), test_video_distribution()"""
    
    # Commands content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = commands_text
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_file_structure_slide(prs):
    """Fayl struktura slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📁 Loyiha Fayl Strukturasi"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # File structure
    file_structure = """📁 **tgbotmuvofiqiyat/** (Asosiy papka)
🔹 **app.py** - Asosiy dastur, bot ishga tushirish
🔹 **loader.py** - Bot va dispatcher yaratish
🔹 **config.py** - Konfiguratsiya va tokenlar
🔹 **db.py** - Database funksiyalari va so'rovlar

📁 **handlers/** (Xabarlarni qayta ishlash)
🔹 **__init__.py** - Handlers import qilish
🔹 **users/** - Foydalanuvchi xabarlari
   📁 **__init__.py** - Users handlers import
   📁 **group_video_commands.py** - Video buyruqlari (asosiy)
   📁 **security.py** - Xavfsizlik va ro'yxatdan o'tish
   📁 **video_scheduler.py** - Video rejalashtirish
   📁 **admin_image_sender.py** - Admin rasm yuborish

📁 **utils/** (Yordamchi funksiyalar)
🔹 **misc/** - Turli xil yordamchi funksiyalar
   📁 **set_bot_commands.py** - Bot buyruqlarini sozlash

📁 **requirements.txt** - Kerakli kutubxonalar
📁 **README.md** - Loyiha haqida ma'lumot"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = file_structure
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("📁") or paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_technical_functions_slide(prs):
    """Texnik funksiyalar slide yaratish"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🛠️ Texnik Funksiyalar va API"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Technical functions
    tech_functions = """🔹 **Database Funksiyalari:**
   • get_all_groups_with_settings() - Barcha guruh sozlamalari
   • update_group_settings() - Guruh sozlamalarini yangilash
   • get_video_keyboard_from_db() - Video klaviaturasini olish
   • check_user_exists() - Foydalanuvchi mavjudligini tekshirish

🔹 **Scheduler Funksiyalari:**
   • schedule_job_with_immediate_check() - Vazifani rejalashtirish
   • send_group_video_new() - Guruhga video yuborish
   • process_video_sending() - Video yuborish jarayoni

🔹 **Security Funksiyalari:**
   • is_admin() - Admin huquqini tekshirish
   • is_super_admin() - Super-admin huquqini tekshirish
   • check_group_whitelist() - Guruh ruxsatini tekshirish

🔹 **Video Handler Funksiyalari:**
   • send_video() - Videoni yuborish
   • send_video_to_group() - Guruhga maxsus video
   • update_video_progress() - Video progress yangilash

🔹 **Keyboard Funksiyalari:**
   • get_project_keyboard() - Loyiha tanlash klaviaturasi
   • get_time_selection_keyboard() - Vaqt tanlash klaviaturasi
   • get_group_selection_keyboard() - Guruh tanlash klaviaturasi"""
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = tech_functions
    
    # Content styling
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        if paragraph.text.startswith("🔹"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def main():
    """Asosiy funksiya"""
    print("🎬 Centris Bot PowerPoint Presentation V4 yaratilmoqda...")
    print("📊 Diagrammalarsiz, funksiyalar va struktura bilan")
    
    # Yangi presentation yaratish
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "🎬 CENTRIS TOWERS & GOLDEN LAKE",
        "Telegram Bot Funksiyalari va Strukturasi\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "📋 Loyiha Haqida",
        """• Centris Towers va Golden Lake loyihalari uchun Telegram bot
• Avtomatik video tarqatish tizimi
• Admin panel va foydalanuvchi interfeysi
• PostgreSQL ma'lumotlar bazasi
• APScheduler bilan vaqt rejalashtirish
• Xavfsizlik tizimi va whitelist
• Professional arxitektura va kod sifat
• Barcha funksiyalar va struktura to'liq"""
    )
    
    # Slide 3: File Structure
    create_file_structure_slide(prs)
    
    # Slide 4: Database Structure
    create_database_structure_slide(prs)
    
    # Slide 5: Bot Commands Detailed
    create_bot_commands_detailed_slide(prs)
    
    # Slide 6: Technical Functions
    create_technical_functions_slide(prs)
    
    # Slide 7: Security System
    create_content_slide(
        prs,
        "🔒 Xavfsizlik Tizimi va Funksiyalari",
        """• Whitelist tizimi - check_group_whitelist()
• Admin roli - is_admin() funksiyasi
• Super-admin roli - is_super_admin() funksiyasi
• Group verification - guruh mavjudligini tekshirish
• Rate limiting - haddan tashqari so'rovlarni cheklash
• Input validation - kiritilgan ma'lumotlarni tekshirish
• Xavfsizlik choralari to'liq amalga oshirilgan
• Barcha funksiyalar xavfsizlik bilan himoyalangan"""
    )
    
    # Slide 8: Video System Functions
    create_functions_slide(
        prs,
        "🎬 Video Tizimi Funksiyalari",
        """🔹 **Video Yuborish:**
   • send_video() - Asosiy video yuborish
   • send_video_to_group() - Guruhga maxsus video
   • send_all_planned_videos() - Barcha rejalashtirilgan

🔹 **Video Boshqaruvi:**
   • update_video_progress() - Progress yangilash
   • get_video_keyboard_from_db() - Video klaviaturasi
   • process_video_sending() - Yuborish jarayoni

🔹 **Vaqt Rejalashtirish:**
   • schedule_job_with_immediate_check() - Vazifa rejalashtirish
   • send_group_video_new() - Guruhga video
   • APScheduler integratsiyasi"""
    )
    
    # Slide 9: Database Functions
    create_functions_slide(
        prs,
        "🗄️ Database Funksiyalari va So'rovlar",
        """🔹 **Guruh Sozlamalari:**
   • get_all_groups_with_settings() - Barcha guruhlar
   • update_group_settings() - Sozlamalarni yangilash
   • check_group_exists() - Guruh mavjudligi

🔹 **Foydalanuvchi Boshqaruvi:**
   • check_user_exists() - Foydalanuvchi mavjudligi
   • register_user() - Yangi foydalanuvchi
   • update_user_subscription() - Obuna yangilash

🔹 **Video Ma'lumotlari:**
   • get_videos_by_season() - Mavsum bo'yicha videolar
   • get_video_by_id() - ID bo'yicha video
   • update_video_progress() - Progress yangilash

🔹 **PostgreSQL Integratsiyasi:**
   • psycopg2 connector
   • Connection pooling
   • Error handling
   • Transaction management"""
    )
    
    # Slide 10: Admin Panel Functions
    create_functions_slide(
        prs,
        "⚙️ Admin Panel Funksiyalari",
        """🔹 **Guruh Boshqaruvi:**
   • set_group_video() - Video sozlamalari
   • show_group_video_settings() - Sozlamalarni ko'rish
   • update_group_names() - Guruh nomlarini yangilash

🔹 **Video Boshqaruvi:**
   • send_specific_video() - Maxsus video yuborish
   • update_video_progress() - Progress yangilash
   • test_send_video_all_groups() - Test video

🔹 **Sistem Boshqaruvi:**
   • is_admin() - Admin huquqini tekshirish
   • is_super_admin() - Super-admin huquqini tekshirish
   • check_group_whitelist() - Guruh ruxsatini tekshirish

🔹 **Xavfsizlik:**
   • Admin buyruqlarini himoya qilish
   • Guruh ruxsatini tekshirish
   • Foydalanuvchi huquqlarini boshqarish"""
    )
    
    # Slide 11: User Experience Functions
    create_functions_slide(
        prs,
        "👥 Foydalanuvchi Tajribasi Funksiyalari",
        """🔹 **Interfeys:**
   • get_project_keyboard() - Loyiha tanlash
   • get_time_selection_keyboard() - Vaqt tanlash
   • get_group_selection_keyboard() - Guruh tanlash
   • get_season_keyboard() - Mavsum tanlash

🔹 **Xabar Yuborish:**
   • send_welcome_message() - Xush kelibsiz xabari
   • send_video_message() - Video xabari
   • send_error_message() - Xatolik xabari
   • send_success_message() - Muvaffaqiyat xabari

🔹 **Foydalanuvchi Boshqaruvi:**
   • register_user() - Ro'yxatdan o'tkazish
   • check_user_exists() - Mavjudligini tekshirish
   • update_user_subscription() - Obunani yangilash

🔹 **Xatoliklarni Bartaraf Etish:**
   • Error handling - Xatoliklarni qayta ishlash
   • Logging - Xatoliklarni qayd qilish
   • User feedback - Foydalanuvchiga ma'lumot"""
    )
    
    # Slide 12: Future Development Functions
    create_content_slide(
        prs,
        "🚀 Kelajak Rivojlantirish Funksiyalari",
        """• Yangi loyihalar qo'shish - add_new_project()
• Video analytics va statistika - get_video_statistics()
• Avtomatik content generation - generate_content()
• AI-powered video recommendations - recommend_videos()
• Mobile app development - mobile_api_integration()
• Integration with CRM systems - crm_integration()
• Multi-language support - multi_language_support()
• Advanced admin dashboard - advanced_admin_panel()
• Cloud deployment va scaling - cloud_deployment()
• Performance monitoring - monitor_performance()"""
    )
    
    # Slide 13: Code Quality Functions
    create_functions_slide(
        prs,
        "✨ Kod Sifati va Xatoliklarni Bartaraf Etish",
        """🔹 **Xatoliklarni Bartaraf Etish:**
   • parse_mode to'g'ri qo'llash
   • Database so'rovlarni optimizatsiya qilish
   • Xatolik xabarlarini qayta ishlash
   • Logging va monitoring

🔹 **Kod Sifati:**
   • PEP 8 standartlariga rioya qilish
   • Docstring va commentlar
   • Error handling va validation
   • Testing va debugging

🔹 **Performance:**
   • Database connection pooling
   • Async/await to'g'ri ishlatish
   • Memory management
   • Response time optimization

🔹 **Security:**
   • Input validation
   • SQL injection himoyasi
   • Rate limiting
   • Access control"""
    )
    
    # Slide 14: Conclusion
    create_content_slide(
        prs,
        "✅ Xulosa va Natijalar",
        """• Centris Bot muvaffaqiyatli ishlayapti
• Barcha xatolar bartaraf etildi
• Video tarqatish tizimi to'liq ishlayapti
• Admin panel qulay va funksional
• Xavfsizlik tizimi ishonchli
• Kod sifatli va tuzilgan
• Professional arxitektura
• Barcha funksiyalar va struktura to'liq
• Kelajakda rivojlantirish imkoniyatlari keng
• Mijozlar mamnun"""
    )
    
    # Slide 15: Thank You
    create_title_slide(
        prs,
        "🎉 Rahmat!",
        "Centris Bot haqida so'rashganingiz uchun!\n\nBarcha funksiyalar va struktura ko'rsatildi!\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Faylni saqlash
    filename = "prezentatsiya_v4_funksiyalar.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint presentation V4 muvaffaqiyatli yaratildi: {filename}")
    print(f"📊 Jami slaydlar soni: {len(prs.slides)}")
    print(f"🔹 Barcha funksiyalar va struktura ko'rsatildi")
    print(f"📱 Bot buyruqlari haqida batafsil ma'lumot")
    print(f"🇺🇿 O'zbek tilida tayyorlandi")
    print(f"💡 Diagrammalarsiz, funksiyalar bilan professional ko'rinish")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎬 {filename} fayli muvaffaqiyatli yaratildi!")
        print("📁 Fayl joriy papkada joylashgan")
        print("🔹 Barcha funksiyalar va struktura to'liq")
        print("💡 Professional ko'rinish diagrammalarsiz")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
        print("💡 python-pptx kutubxonasi o'rnatilganligini tekshiring")
