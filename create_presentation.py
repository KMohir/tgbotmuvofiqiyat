#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Centris Bot PowerPoint Presentation Creator
Uzbek tilida diagrammalar bilan
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    """Title slide yaratish"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Styling
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_content_slide(prs, title, content, diagram_type=None):
    """Content slide yaratish"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    
    # Content
    if diagram_type:
        content_shape.text = f"{content}\n\n📊 {diagram_type}"
    else:
        content_shape.text = content
    
    # Content styling
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
    
    return slide

def create_diagram_slide(prs, title, diagram_text, description):
    """Diagramma slide yaratish"""
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Diagram
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    diagram_frame = diagram_box.text_frame
    diagram_frame.text = diagram_text
    diagram_frame.paragraphs[0].font.size = Pt(14)
    diagram_frame.paragraphs[0].font.name = 'Courier New'
    diagram_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Asosiy funksiya"""
    print("🎬 Centris Bot PowerPoint Presentation yaratilmoqda...")
    
    # Yangi presentation yaratish
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "🎬 CENTRIS TOWERS & GOLDEN LAKE",
        "Telegram Bot Arxitekturasi va Ishlash Prinsipi\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "📋 Loyiha Haqida",
        """• Centris Towers va Golden Lake loyihalari uchun Telegram bot
• Avtomatik video tarqatish tizimi
• Admin panel va foydalanuvchi interfeysi
• PostgreSQL ma'lumotlar bazasi
• APScheduler bilan vaqt rejalashtirish""",
        "Loyiha tuzilishi"
    )
    
    # Slide 3: Architecture Diagram
    create_diagram_slide(
        prs,
        "🏗️ Arxitektura Diagrammasi",
        """┌─────────────────┐    ┌─────────────────┐    ┌─────────────────┐
│   Telegram API  │    │   Centris Bot   │    │  PostgreSQL DB  │
│                 │◄──►│                 │◄──►│                 │
│  User Messages  │    │  Video Handler  │    │  User Data      │
│  Admin Commands │    │  Scheduler      │    │  Video Settings │
└─────────────────┘    │  Security       │    │  Group Settings │
                       └─────────────────┘    └─────────────────┘
                                │
                                ▼
                       ┌─────────────────┐
                       │  Video Storage  │
                       │  (Telegram)     │
                       └─────────────────┘""",
        "Bot qanday ishlayotgani va qaysi komponentlar bilan bog'langan"
    )
    
    # Slide 4: Database Schema
    create_diagram_slide(
        prs,
        "🗄️ Ma'lumotlar Bazasi Tuzilishi",
        """┌─────────────────────┐    ┌─────────────────────┐
│      users           │    │   group_whitelist   │
│  - id (PK)          │    │  - group_id (PK)    │
│  - username         │    │  - group_name       │
│  - subscription     │    │  - added_date       │
└─────────────────────┘    └─────────────────────┘
           │                           │
           │                           │
           ▼                           ▼
┌─────────────────────┐    ┌─────────────────────┐
│  group_video_       │    │      admins         │
│  settings           │    │  - user_id (PK)     │
│  - group_id (PK)    │    │  - role             │
│  - centris_enabled  │    │  - permissions      │
│  - golden_enabled   │    └─────────────────────┘
│  - send_times       │
└─────────────────────┘
           │
           ▼
┌─────────────────────┐    ┌─────────────────────┐
│      seasons        │    │      videos         │
│  - id (PK)          │    │  - id (PK)          │
│  - name             │    │  - season_id (FK)   │
│  - project          │    │  - url              │
└─────────────────────┘    └─────────────────────┘""",
        "Ma'lumotlar bazasida qanday jadvallar va ular o'rtasidagi bog'lanishlar"
    )
    
    # Slide 5: Video Distribution Flow
    create_diagram_slide(
        prs,
        "🎬 Video Tarqatish Jarayoni",
        """┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   Admin     │    │  Scheduler │    │    Bot     │
│   Panel     │    │            │    │            │
└─────┬───────┘    └─────┬──────┘    └─────┬──────┘
      │                  │                 │
      │ Set Time         │                 │
      │ 08:00, 20:00    │                 │
      ├─────────────────►│                 │
      │                  │                 │
      │                  │ Time Check      │
      │                  │ 08:00 = True    │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │                 │ Send Video
      │                  │                 │ to Groups
      │                  │                 ├─────────►
      │                  │                 │
      │                  │                 │ Update DB
      │                  │                 │ Progress""",
        "Video qanday vaqtda va qanday tartibda guruhlarga yuboriladi"
    )
    
    # Slide 6: User Registration Flow
    create_diagram_slide(
        prs,
        "👤 Foydalanuvchi Ro'yxatdan O'tishi",
        """┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   User      │    │    Bot     │    │     DB      │
│  /start     │    │            │    │             │
└─────┬───────┘    └─────┬──────┘    └─────┬───────┘
      │                  │                 │
      │ Send /start      │                 │
      ├─────────────────►│                 │
      │                  │                 │
      │                  │ Check User      │
      │                  │ in Database    │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │ User Not Found │
      │                  │◄────────────────┤
      │                  │                 │
      │                  │ Create New User│
      │                  ├─────────────────►│
      │                  │                 │
      │ Welcome Message  │                 │
      │◄────────────────┤                 │""",
        "Yangi foydalanuvchi qanday ro'yxatdan o'tadi va bot qanday javob beradi"
    )
    
    # Slide 7: Admin Panel Workflow
    create_diagram_slide(
        prs,
        "⚙️ Admin Panel Ishlash Tartibi",
        """┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   Admin     │    │    Bot     │    │     DB      │
│  Commands   │    │            │    │             │
└─────┬───────┘    └─────┬──────┘    └─────┬───────┘
      │                  │                 │
      │ /set_group_video │                 │
      ├─────────────────►│                 │
      │                  │                 │
      │                  │ Check Admin     │
      │                  │ Permissions     │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │ Admin Verified  │
      │                  │◄────────────────┤
      │                  │                 │
      │                  │ Update Group    │
      │                  │ Settings        │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │ Update Scheduler│
      │                  ├─────────────────►│
      │                  │                 │
      │ Success Message  │                 │
      │◄────────────────┤                 │""",
        "Admin qanday buyruqlar orqali bot sozlamalarini o'zgartiradi"
    )
    
    # Slide 8: Security System
    create_content_slide(
        prs,
        "🔒 Xavfsizlik Tizimi",
        """• Whitelist tizimi - faqat ruxsat berilgan guruhlar
• Admin roli - faqat admin buyruqlarini bajaradi
• Super-admin roli - barcha huquqlarga ega
• Group verification - guruh mavjudligini tekshirish
• Rate limiting - haddan tashqari so'rovlarni cheklash
• Input validation - kiritilgan ma'lumotlarni tekshirish""",
        "Xavfsizlik choralari"
    )
    
    # Slide 9: Technical Stack
    create_content_slide(
        prs,
        "🛠️ Texnik Stack",
        """• Python 3.13 - asosiy dasturlash tili
• aiogram 2.x - Telegram Bot API kutubxonasi
• PostgreSQL - ma'lumotlar bazasi
• APScheduler - vaqt rejalashtirish
• psycopg2 - PostgreSQL connector
• pytz - vaqt zonasi boshqaruvi
• aiohttp - HTTP client kutubxonasi""",
        "Qanday texnologiyalar ishlatilgan"
    )
    
    # Slide 10: Commands Overview
    create_content_slide(
        prs,
        "📱 Bot Buyruqlari",
        """• /start - botni ishga tushirish
• /centris_towers - Centris videolarini ko'rish
• /golden_lake - Golden Lake videolarini ko'rish
• /set_group_video - guruh uchun video sozlamalari
• /show_group_video_settings - guruh sozlamalarini ko'rish
• /update_video_progress - video progress yangilash
• /send_specific_video - maxsus video yuborish""",
        "Foydalanuvchilar va adminlar uchun mavjud buyruqlar"
    )
    
    # Slide 11: Video Scheduling Algorithm
    create_diagram_slide(
        prs,
        "⏰ Video Rejalashtirish Algoritmi",
        """┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   Time      │    │ Scheduler │    │   Groups    │
│   Check     │    │           │    │             │
└─────┬───────┘    └─────┬──────┘    └─────┬───────┘
      │                  │                 │
      │ Current Time     │                 │
      │ 08:00 = True    │                 │
      ├─────────────────►│                 │
      │                  │                 │
      │                  │ Get All Groups │
      │                  │ with Settings  │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │ Filter by Time │
      │                  │ 08:00 Groups   │
      │                  │◄────────────────┤
      │                  │                 │
      │                  │ Send Videos    │
      │                  │ to Each Group  │
      │                  ├─────────────────►│
      │                  │                 │
      │                  │ Update Progress │
      │                  │ in Database    │""",
        "Video qanday vaqtda va qaysi guruhlarga yuboriladi"
    )
    
    # Slide 12: Future Plans
    create_content_slide(
        prs,
        "🚀 Kelajak Rejalari",
        """• Yangi loyihalar qo'shish
• Video analytics va statistika
• Avtomatik content generation
• AI-powered video recommendations
• Mobile app development
• Integration with CRM systems
• Multi-language support
• Advanced admin dashboard""",
        "Loyihaning rivojlanish yo'nalishlari"
    )
    
    # Slide 13: Conclusion
    create_content_slide(
        prs,
        "✅ Xulosa",
        """• Centris Bot muvaffaqiyatli ishlayapti
• Barcha xatolar bartaraf etildi
• Video tarqatish tizimi to'liq ishlayapti
• Admin panel qulay va funksional
• Xavfsizlik tizimi ishonchli
• Kod sifatli va tuzilgan
• Kelajakda rivojlantirish imkoniyatlari keng""",
        "Loyihaning hozirgi holati va natijalari"
    )
    
    # Slide 14: Thank You
    create_title_slide(
        prs,
        "🎉 Rahmat!",
        "Centris Bot haqida so'rashganingiz uchun!\n\nSavollar bo'lsa, so'rang!\n\nTayyorladi: AI Assistant\nSana: 31.08.2025"
    )
    
    # Faylni saqlash
    filename = "prezentatsiya.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint presentation muvaffaqiyatli yaratildi: {filename}")
    print(f"📊 Jami slaydlar soni: {len(prs.slides)}")
    print(f"🎯 Diagrammalar qo'shildi")
    print(f"🇺🇿 O'zbek tilida tayyorlandi")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎬 {filename} fayli muvaffaqiyatli yaratildi!")
        print("📁 Fayl joriy papkada joylashgan")
    except Exception as e:
        print(f"❌ Xatolik yuz berdi: {e}")
        print("💡 python-pptx kutubxonasi o'rnatilganligini tekshiring")
